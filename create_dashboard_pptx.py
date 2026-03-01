#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Color palette
colors = {
    'primary': RGBColor(0x25, 0x63, 0xEB),      # Blue
    'secondary': RGBColor(0x7C, 0x3A, 0xED),    # Purple
    'success': RGBColor(0x16, 0xA3, 0x4A),      # Green
    'warning': RGBColor(0xEA, 0x58, 0x0C),      # Orange
    'dark': RGBColor(0x11, 0x18, 0x27),         # Dark gray
    'lightText': RGBColor(0x6B, 0x72, 0x80),    # Light gray
    'bg': RGBColor(0xF8, 0xFA, 0xFC),          # Off-white
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'card': RGBColor(0xF9, 0xFA, 0xFB),        # Card bg
    'border': RGBColor(0xE5, 0xE7, 0xEB)       # Border
}

def add_metric_card(slide, left, top, label, value, subtext=''):
    """Add a metric card"""
    # Background
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(left), Inches(top), Inches(1.8), Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors['white']
    shape.line.color.rgb = colors['border']

    # Label
    text_frame = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.08), Inches(1.6), Inches(0.3)
    ).text_frame
    p = text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = colors['lightText']
    p.font.bold = True

    # Value
    text_frame = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.35), Inches(1.6), Inches(0.45)
    ).text_frame
    p = text_frame.paragraphs[0]
    p.text = value
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = colors['dark']

    # Subtext
    if subtext:
        text_frame = slide.shapes.add_textbox(
            Inches(left + 0.1), Inches(top + 0.8), Inches(1.6), Inches(0.25)
        ).text_frame
        p = text_frame.paragraphs[0]
        p.text = subtext
        p.font.size = Pt(9)
        p.font.color.rgb = colors['lightText']
        p.font.italic = True

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = colors['primary']

# Title
text_frame = slide1.shapes.add_textbox(
    Inches(0.5), Inches(1.5), Inches(9), Inches(1)
).text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "Ad Auction Simulator"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = colors['white']
p.alignment = PP_ALIGN.CENTER

# Subtitle
text_frame = slide1.shapes.add_textbox(
    Inches(0.5), Inches(2.6), Inches(9), Inches(0.6)
).text_frame
p = text_frame.paragraphs[0]
p.text = "GSP/VCG Auction Engine · Recommender Model Routing · LLM-Powered Analysis"
p.font.size = Pt(18)
p.font.color.rgb = colors['white']
p.alignment = PP_ALIGN.CENTER

# Tagline
text_frame = slide1.shapes.add_textbox(
    Inches(0.5), Inches(4.0), Inches(9), Inches(0.5)
).text_frame
p = text_frame.paragraphs[0]
p.text = "Interactive Dashboard Demo"
p.font.size = Pt(16)
p.font.italic = True
p.font.color.rgb = colors['white']
p.alignment = PP_ALIGN.CENTER

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: Auction Dashboard
# ═══════════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide2.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = colors['bg']

# Title
text_frame = slide2.shapes.add_textbox(
    Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
).text_frame
p = text_frame.paragraphs[0]
p.text = "Auction Dashboard"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = colors['dark']

# Tabs
tabs = ["Auction Dashboard", "Segment Explorer", "What-If Analysis"]
for i, tab in enumerate(tabs):
    is_active = i == 0
    shape = slide2.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5 + i * 2.5), Inches(0.85), Inches(2.3), Inches(0.35)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors['white'] if is_active else colors['bg']
    shape.line.color.rgb = colors['primary'] if is_active else colors['border']

    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = tab
    p.font.size = Pt(10)
    p.font.bold = is_active
    p.font.color.rgb = colors['primary'] if is_active else colors['lightText']
    p.alignment = PP_ALIGN.CENTER

# Metric Cards
add_metric_card(slide2, 0.5, 1.4, "Total Revenue", "$2,847.32", "Per 1K impr.")
add_metric_card(slide2, 2.6, 1.4, "Avg RPM", "$4.87", "Rev per 1K impr")
add_metric_card(slide2, 4.7, 1.4, "Total Clicks", "14,287", "Across segments")
add_metric_card(slide2, 6.8, 1.4, "Advertisers", "80", "8 verticals")

# Table: Auction Winners
text_frame = slide2.shapes.add_textbox(
    Inches(0.5), Inches(2.8), Inches(9), Inches(0.3)
).text_frame
p = text_frame.paragraphs[0]
p.text = "Auction Winners — Young Tech Enthusiasts (GSP Mechanism)"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = colors['dark']

# Simple table visualization
table_data = [
    ["Slot", "Advertiser", "Vertical", "Bid", "Quality", "Eff. Bid", "CPC", "pCTR"],
    ["1", "Finance Adv. 23", "Finance", "$9.54", "0.891", "8.504", "$6.21", "4.21%"],
    ["2", "E-Comm Adv. 5", "E-Commerce", "$7.82", "0.764", "5.974", "$3.89", "3.87%"],
    ["3", "SaaS Adv. 12", "SaaS", "$6.45", "0.701", "4.522", "$2.45", "3.45%"],
    ["4", "Gaming Adv. 8", "Gaming", "$5.12", "0.623", "3.190", "$1.87", "3.12%"]
]

y_pos = 3.15
for row_idx, row in enumerate(table_data):
    x_pos = 0.5
    is_header = row_idx == 0
    row_bg = colors['card'] if row_idx % 2 == 1 else colors['white']

    for col_idx, cell in enumerate(row):
        # Cell background
        shape = slide2.shapes.add_shape(1, Inches(x_pos), Inches(y_pos), Inches(1.1), Inches(0.32))
        shape.fill.solid()
        shape.fill.fore_color.rgb = row_bg
        shape.line.color.rgb = colors['border']
        shape.line.width = Pt(0.5)

        # Cell text
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = cell
        p.font.size = Pt(8)
        p.font.bold = is_header
        p.font.color.rgb = colors['dark']
        p.alignment = PP_ALIGN.CENTER

        x_pos += 1.1

    y_pos += 0.33

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: Reserve Price Analysis
# ═══════════════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide3.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = colors['bg']

# Title
text_frame = slide3.shapes.add_textbox(
    Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
).text_frame
p = text_frame.paragraphs[0]
p.text = "Revenue vs Reserve Price Analysis"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = colors['dark']

# Chart area
chart_shape = slide3.shapes.add_shape(
    1,  # Rectangle
    Inches(0.5), Inches(1.0), Inches(6.5), Inches(3.8)
)
chart_shape.fill.solid()
chart_shape.fill.fore_color.rgb = colors['white']
chart_shape.line.color.rgb = colors['border']

# Bar chart (simplified)
bar_heights = [2.8, 2.5, 2.2, 2.0, 2.3, 3.0, 3.5, 3.2, 2.5, 1.8]
for i, height in enumerate(bar_heights):
    bar = slide3.shapes.add_shape(
        1,
        Inches(0.7 + i * 0.55), Inches(4.2 - height * 0.4),
        Inches(0.4), Inches(height * 0.4)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors['primary']
    bar.line.color.rgb = colors['primary']

# Optimal point
optimal = slide3.shapes.add_shape(
    3,  # Oval
    Inches(2.9), Inches(3.8), Inches(0.2), Inches(0.2)
)
optimal.fill.solid()
optimal.fill.fore_color.rgb = colors['warning']
optimal.line.color.rgb = colors['dark']

text_frame = slide3.shapes.add_textbox(
    Inches(2.3), Inches(3.5), Inches(1.5), Inches(0.6)
).text_frame
p = text_frame.paragraphs[0]
p.text = "OPTIMAL\n$2.00 → $3,421"
p.font.size = Pt(9)
p.font.bold = True
p.font.color.rgb = colors['warning']
p.alignment = PP_ALIGN.CENTER

# Axis labels
text_frame = slide3.shapes.add_textbox(
    Inches(0.7), Inches(4.5), Inches(5.0), Inches(0.3)
).text_frame
p = text_frame.paragraphs[0]
p.text = "$0    $1    $2    $3    $4    $5"
p.font.size = Pt(10)
p.font.color.rgb = colors['lightText']

# Key Insights panel
text_frame = slide3.shapes.add_textbox(
    Inches(7.2), Inches(1.0), Inches(2.3), Inches(0.35)
).text_frame
p = text_frame.paragraphs[0]
p.text = "Key Insights"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = colors['dark']

insights = [
    "• Sweet spot: $2.00",
    "• Max revenue: $3,421",
    "• Fill rate: 58%",
    "• vs $0.50: +20% rev",
    "• vs $5.00: -44% rev"
]

y_pos = 1.4
for insight in insights:
    text_frame = slide3.shapes.add_textbox(
        Inches(7.2), Inches(y_pos), Inches(2.3), Inches(0.3)
    ).text_frame
    p = text_frame.paragraphs[0]
    p.text = insight
    p.font.size = Pt(10)
    p.font.color.rgb = colors['dark']
    y_pos += 0.35

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Segment Explorer
# ═══════════════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide4.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = colors['bg']

# Title
text_frame = slide4.shapes.add_textbox(
    Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
).text_frame
p = text_frame.paragraphs[0]
p.text = "Segment Explorer & Model Routing"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = colors['dark']

# Segment buttons
segments = ["Young Tech", "Suburban P.", "Luxury", "College", "Biz Prof.", "Fitness", "Gamers", "Retirees"]
btn_x = 0.5
for i, seg in enumerate(segments):
    is_active = i == 0
    btn = slide4.shapes.add_shape(
        1,
        Inches(btn_x), Inches(1.0), Inches(1.08), Inches(0.3)
    )
    btn.fill.solid()
    btn.fill.fore_color.rgb = colors['primary'] if is_active else colors['white']
    btn.line.color.rgb = colors['primary'] if is_active else colors['border']

    text_frame = btn.text_frame
    p = text_frame.paragraphs[0]
    p.text = seg
    p.font.size = Pt(8)
    p.font.bold = is_active
    p.font.color.rgb = colors['white'] if is_active else colors['lightText']
    p.alignment = PP_ALIGN.CENTER

    btn_x += 1.12

# Segment info cards
add_metric_card(slide4, 0.5, 1.5, "Segment Size", "2.4M", "")
add_metric_card(slide4, 2.6, 1.5, "Avg CTR", "4.2%", "")
add_metric_card(slide4, 4.7, 1.5, "Avg CVR", "1.8%", "")
add_metric_card(slide4, 6.8, 1.5, "Best Model", "DLRM", "")

# Model table
text_frame = slide4.shapes.add_textbox(
    Inches(0.5), Inches(2.95), Inches(6.0), Inches(0.3)
).text_frame
p = text_frame.paragraphs[0]
p.text = "Model Performance Comparison"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = colors['dark']

model_data = [
    ["Model", "CTR Lift", "Rev Lift", "Latency"],
    ["Two-Tower", "0.785", "0.746", "5ms"],
    ["GBDT", "0.812", "0.633", "12ms"],
    ["DLRM", "0.860", "0.731", "24ms"],
    ["Bandit", "0.755", "0.619", "8ms"]
]

y_pos = 3.35
col_widths = [1.5, 1.5, 1.5, 0.9]
for row_idx, row in enumerate(model_data):
    x_pos = 0.5
    is_header = row_idx == 0
    is_recommended = row_idx == 3
    row_bg = colors['white'] if is_header else (colors['primary'] if is_recommended else colors['card'])
    row_text_color = colors['white'] if is_recommended else colors['dark']

    for col_idx, cell in enumerate(row):
        cell_shape = slide4.shapes.add_shape(
            1,
            Inches(x_pos), Inches(y_pos), Inches(col_widths[col_idx]), Inches(0.32)
        )
        cell_shape.fill.solid()
        cell_shape.fill.fore_color.rgb = row_bg
        cell_shape.line.color.rgb = colors['border']
        cell_shape.line.width = Pt(0.5)

        text_frame = cell_shape.text_frame
        p = text_frame.paragraphs[0]
        p.text = cell
        p.font.size = Pt(9)
        p.font.bold = is_header or is_recommended
        p.font.color.rgb = row_text_color
        p.alignment = PP_ALIGN.CENTER

        x_pos += col_widths[col_idx]

    y_pos += 0.33

# Recommendation banner
banner = slide4.shapes.add_shape(
    1,
    Inches(0.5), Inches(4.95), Inches(6.0), Inches(0.5)
)
banner.fill.solid()
banner.fill.fore_color.rgb = RGBColor(0x25, 0x63, 0xEB)
banner.fill.transparency = 0.9
banner.line.color.rgb = colors['primary']

text_frame = banner.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "✓ RECOMMENDED: DLRM Deep Model — Warm segment with high data density. Revenue lift 0.731 justifies compute cost."
p.font.size = Pt(9)
p.font.bold = True
p.font.color.rgb = colors['dark']

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: What-If Chat
# ═══════════════════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide5.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = colors['bg']

# Title
text_frame = slide5.shapes.add_textbox(
    Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
).text_frame
p = text_frame.paragraphs[0]
p.text = "What-If Analysis Chat"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = colors['dark']

# Chat box
chat_box = slide5.shapes.add_shape(
    1,
    Inches(0.5), Inches(1.0), Inches(9), Inches(3.8)
)
chat_box.fill.solid()
chat_box.fill.fore_color.rgb = colors['white']
chat_box.line.color.rgb = colors['border']

# Assistant message
assistant_msg = slide5.shapes.add_shape(
    1,
    Inches(0.7), Inches(1.15), Inches(8.6), Inches(0.8)
)
assistant_msg.fill.solid()
assistant_msg.fill.fore_color.rgb = colors['card']
assistant_msg.line.type = None

text_frame = assistant_msg.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "Welcome to the What-If Analyzer. I can run simulations and analyze the impact on revenue, advertiser surplus, and segment performance."
p.font.size = Pt(9)
p.font.italic = True
p.font.color.rgb = colors['dark']

# User message
user_msg = slide5.shapes.add_shape(
    1,
    Inches(3.5), Inches(2.15), Inches(5.8), Inches(0.65)
)
user_msg.fill.solid()
user_msg.fill.fore_color.rgb = RGBColor(0x25, 0x63, 0xEB)
user_msg.fill.transparency = 0.9
user_msg.line.color.rgb = colors['primary']
user_msg.line.width = Pt(0.5)

text_frame = user_msg.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "What if we raise reserve prices to $2.50?"
p.font.size = Pt(9)
p.font.bold = True
p.font.color.rgb = colors['dark']

# Response message
response_box = slide5.shapes.add_shape(
    1,
    Inches(0.7), Inches(2.95), Inches(8.6), Inches(1.65)
)
response_box.fill.solid()
response_box.fill.fore_color.rgb = colors['card']
response_box.line.type = None

text_frame = response_box.text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "Reserve Price Impact Analysis\n\nRevenue: $2,847 → $3,156 (+10.8%)\n\nHigher reserve prices filter low-quality bids, increasing average CPC. Watch for fill rate degradation in smaller segments."
p.font.size = Pt(9)
p.font.color.rgb = colors['dark']

# Input field
input_box = slide5.shapes.add_shape(
    1,
    Inches(0.5), Inches(5.0), Inches(8.5), Inches(0.4)
)
input_box.fill.solid()
input_box.fill.fore_color.rgb = colors['white']
input_box.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)

text_frame = input_box.text_frame
p = text_frame.paragraphs[0]
p.text = "Ask a what-if question about the auction system..."
p.font.size = Pt(10)
p.font.italic = True
p.font.color.rgb = colors['lightText']

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Key Features
# ═══════════════════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide6.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = colors['bg']

# Title
text_frame = slide6.shapes.add_textbox(
    Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
).text_frame
p = text_frame.paragraphs[0]
p.text = "Key Features & Capabilities"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = colors['dark']

features = [
    ("📊", "Auction Engine", "GSP & VCG mechanisms"),
    ("🔄", "Model Routing", "4 recommender models"),
    ("💬", "What-If Analysis", "Claude-powered queries"),
    ("📈", "Interactive Charts", "Real-time metrics"),
    ("🎯", "Sensitivity Analysis", "Parameter impact"),
    ("🏆", "Simulation Data", "80 synthetic advertisers")
]

feature_y = 1.2
for idx, (icon, title, desc) in enumerate(features):
    col_x = 0.5 if idx < 3 else 5.3
    row_offset = idx % 3
    y_pos = feature_y + row_offset * 1.15

    # Icon circle
    icon_shape = slide6.shapes.add_shape(
        3,  # Oval
        Inches(col_x), Inches(y_pos), Inches(0.45), Inches(0.45)
    )
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = colors['primary']
    icon_shape.fill.transparency = 0.8
    icon_shape.line.type = None

    # Icon text
    text_frame = icon_shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = icon
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER

    # Title
    text_frame = slide6.shapes.add_textbox(
        Inches(col_x + 0.55), Inches(y_pos), Inches(4.2), Inches(0.25)
    ).text_frame
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = colors['dark']

    # Description
    text_frame = slide6.shapes.add_textbox(
        Inches(col_x + 0.55), Inches(y_pos + 0.25), Inches(4.2), Inches(0.2)
    ).text_frame
    p = text_frame.paragraphs[0]
    p.text = desc
    p.font.size = Pt(9)
    p.font.color.rgb = colors['lightText']

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: Closing
# ═══════════════════════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide7.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = colors['secondary']

# Main text
text_frame = slide7.shapes.add_textbox(
    Inches(0.5), Inches(1.8), Inches(9), Inches(0.7)
).text_frame
p = text_frame.paragraphs[0]
p.text = "Ready to Explore?"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = colors['white']
p.alignment = PP_ALIGN.CENTER

# Subtitle
text_frame = slide7.shapes.add_textbox(
    Inches(0.5), Inches(2.6), Inches(9), Inches(0.5)
).text_frame
p = text_frame.paragraphs[0]
p.text = "Run the dashboard locally and start experimenting with auction dynamics"
p.font.size = Pt(15)
p.font.color.rgb = colors['white']
p.alignment = PP_ALIGN.CENTER

# Instructions
text_frame = slide7.shapes.add_textbox(
    Inches(1.5), Inches(3.5), Inches(7), Inches(1.2)
).text_frame
text_frame.word_wrap = True
p = text_frame.paragraphs[0]
p.text = "Backend: uvicorn app.main:app --reload --port 8000\n\nFrontend: npm run dev (http://localhost:3000)"
p.font.size = Pt(11)
p.font.color.rgb = colors['white']
p.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = "/sessions/vigilant-affectionate-archimedes/mnt/outputs/Ad_Auction_Simulator_Dashboard.pptx"
prs.save(output_path)
print(f"✓ Presentation created: Ad_Auction_Simulator_Dashboard.pptx")
